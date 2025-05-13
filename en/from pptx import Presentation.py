from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

# Add Title
title_box = slide.shapes.add_textbox(Inches(2), Inches(0.3), Inches(6), Inches(1))
tf = title_box.text_frame
tf.text = "ZIUA MAMEI"
tf.paragraphs[0].font.size = Pt(48)
tf.paragraphs[0].font.bold = True

# Add Church Name
subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(7), Inches(1))
tf = subtitle_box.text_frame
tf.text = "Biserica Sf. Athanasie cel Mare din Aurora"
tf.paragraphs[0].font.size = Pt(24)

# Add Body Text
body_text = (
    "Vă invită la vânzarea organizată pentru soții, mame, fiice, surori, bunice, prietene.\n"
    "11 și 25 Mai 2025\n"
    "Vino după slujbă, să găsești cadoul potrivit cu ajutorul voluntarelor Organizației de Femei Sf. Ana."
)
body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(2))
tf = body_box.text_frame
for line in body_text.split("\n"):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(20)

# Add Location and Time
info_text = (
    "Unde?\n"
    "În sala socială a bisericii noastre\n"
    "Când?\n"
    "Imediat după slujbă"
)
info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(8), Inches(2))
tf = info_box.text_frame
for line in info_text.split("\n"):
    p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    if line.endswith("?"):
        p.font.bold = True

# Add Contact Information
contact_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
tf = contact_box.text_frame
tf.text = "Informații suplimentare la numărul de telefon: 630-8819529"
tf.paragraphs[0].font.size = Pt(16)

# Save presentation
prs.save('/mnt/data/ziua_mamei.pptx')
